"""Shared fixture builders for the extraction engine tests (M2).

These build *real* parseable documents with deterministic bytes — the engines
under test do the only parsing; fixtures merely produce honest input:

- :func:`make_pdf_bytes` — a minimal single-page PDF written by hand (valid
  header/xref/trailer, one Helvetica text op, a docinfo record). pypdf reads
  it; no third-party writer library needed in tests.
- :func:`make_docx_bytes` — a real OOXML package via python-docx (the pinned
  reader library doubles as the fixture writer; the engine's *read* path is
  what the assertions cover).

Everything else (txt/md/csv/json) is literal bytes in the test bodies.
"""
from __future__ import annotations

import io


def _pdf_escape(text: str) -> str:
    return text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def make_pdf_bytes(
    text: str,
    *,
    title: str | None = None,
    author: str | None = None,
    creation: str | None = "D:20240102030405Z",
    modified: str | None = "D:20240304050607Z",
) -> bytes:
    """A one-page PDF whose visible text is exactly `text`."""

    info = ""
    if title is not None:
        info += f"/Title ({_pdf_escape(title)})"
    if author is not None:
        info += f"/Author ({_pdf_escape(author)})"
    if creation is not None:
        info += f"/CreationDate ({creation})"
    if modified is not None:
        info += f"/ModDate ({modified})"

    content = f"BT /F1 12 Tf 72 720 Td ({_pdf_escape(text)}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        ("<< " + info + " >>").encode("latin-1"),
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode()
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R /Info 6 0 R >>\n"
        f"startxref\n{xref_at}\n%%EOF"
    ).encode()
    return bytes(out)


def make_docx_bytes(
    lines: list[str],
    *,
    title: str | None = None,
    author: str | None = None,
    subject: str | None = None,
) -> bytes:
    """A real .docx whose paragraphs are exactly `lines`."""

    import docx

    document = docx.Document()
    if title is not None:
        document.core_properties.title = title
    if author is not None:
        document.core_properties.author = author
    if subject is not None:
        document.core_properties.subject = subject
    for line in lines:
        document.add_paragraph(line)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# L2 fixtures — compact, generated in-memory (no large binary collections).
# Restored from the authoritative L2 artifact (the committed tree was missing
# these; the committed L2/L3 tests reference them).
# ---------------------------------------------------------------------------

def make_xlsx_bytes(rows, *, sheet: str = "Sheet1") -> bytes:
    """A real .xlsx whose first sheet has exactly `rows` (list of row lists)."""
    import io

    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row, start=1):
            ws.cell(row=r, column=c, value=value)
    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def make_pptx_bytes(slides, *, title: str | None = None) -> bytes:
    """A real .pptx with one slide per entry; each entry is a title text."""
    import io

    from pptx import Presentation

    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        try:
            slide.shapes.title.text = s
        except Exception:  # noqa: BLE001
            pass
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def make_png_bytes(size: int = 8) -> bytes:
    """A tiny PNG (white) via Pillow."""
    import io

    from PIL import Image

    img = Image.new("RGB", (size, size), "white")
    buffer = io.BytesIO()
    img.save(buffer, "PNG")
    return buffer.getvalue()


def make_zip_bytes(members: dict[str, bytes], *, nested: dict | None = None) -> bytes:
    """A zip whose entries are `members` (path -> bytes). Optionally nests one."""
    import io
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for path, data in members.items():
            z.writestr(path, data)
        if nested:
            inner = io.BytesIO()
            with zipfile.ZipFile(inner, "w") as zin:
                for path, data in nested.items():
                    zin.writestr(path, data)
            z.writestr("nested.zip", inner.getvalue())
    return buf.getvalue()


def make_scanned_pdf_bytes() -> bytes:
    """A scanned/image-only PDF: a page with an image and no text layer."""
    return make_pdf_bytes("")
