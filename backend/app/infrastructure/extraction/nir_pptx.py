"""L2 PPTX engine adapter (python-pptx).

Presentation / slide / text region / table / image / shape / speaker-notes
provenance. Slides are preserved as structural units (not flattened).
"""

from __future__ import annotations

import io

from app.application.dtos.nir import NirDocument, NirElement, NirElementType, NirImage
from app.application.ports.nir_parser import NirParseError, NirParser
from app.application.services.extraction_limits import MAX_SLIDES
from app.domain.value_objects.source import MediaKind


class PptxNirParser(NirParser):
    format_name = "pptx"

    def __init__(self) -> None:
        self._engine = "python-pptx"

    def parse(self, data: bytes, *, source_id: str, version: int) -> NirDocument:
        try:
            import pptx
        except ImportError as exc:  # pragma: no cover
            raise NirParseError(f"PPTX engine unavailable: {exc}.") from exc

        try:
            prs = pptx.Presentation(io.BytesIO(data))
        except Exception as exc:  # noqa: BLE001
            raise NirParseError(f"PPTX could not be parsed ({type(exc).__name__}: {exc}).") from exc

        elements: list[NirElement] = []
        images: list[NirImage] = []
        text_parts: list[str] = []
        order = 0
        slide_count = min(len(prs.slides), MAX_SLIDES)

        for s_index in range(slide_count):
            try:
                slide = prs.slides[s_index]
            except Exception:  # noqa: BLE001
                continue
            elements.append(
                NirElement(
                    element_type=NirElementType.SLIDE, order=order,
                    text=f"Slide {s_index + 1}", slide=s_index + 1,
                    extraction_confidence=1.0,
                )
            )
            order += 1
            # text regions from shapes
            for shape in slide.shapes:
                text = ""
                try:
                    if shape.has_text_frame:
                        text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                except Exception:  # noqa: BLE001
                    text = ""
                if text:
                    text_parts.append(text)
                    elements.append(
                        NirElement(
                            element_type=NirElementType.TEXT, order=order,
                            text=text, slide=s_index + 1,
                            extraction_confidence=1.0,
                        )
                    )
                    order += 1
                # tables
                try:
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        rows = [[c.text or "" for c in row.cells] for row in table.rows]
                        elements.append(
                            NirElement(
                                element_type=NirElementType.TABLE, order=order,
                                text="\n".join(" | ".join(r) for r in rows),
                                value={"slide": s_index + 1, "rows": rows},
                                slide=s_index + 1, extraction_confidence=1.0,
                            )
                        )
                        order += 1
                except Exception:  # noqa: BLE001
                    pass
                # images
                try:
                    if getattr(shape, "shape_type", None) is not None and shape.shape_type == 13:
                        image_id = f"pptx-{source_id}-s{s_index + 1}-{shape.shape_id}"
                        bbox = None
                        try:
                            bbox = (shape.left, shape.top, shape.left + shape.width,
                                    shape.top + shape.height)
                        except Exception:  # noqa: BLE001
                            bbox = None
                        images.append(
                            NirImage(
                                image_id=image_id, slide=s_index + 1, bbox=bbox,
                                region={"shape_id": shape.shape_id},
                                media_type="image",
                            )
                        )
                        elements.append(
                            NirElement(
                                element_type=NirElementType.IMAGE, order=order,
                                text="", value={"image_id": image_id},
                                slide=s_index + 1, bbox=bbox, extraction_confidence=1.0,
                            )
                        )
                        order += 1
                except Exception:  # noqa: BLE001
                    pass
            # speaker notes
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes and notes.strip():
                        elements.append(
                            NirElement(
                                element_type=NirElementType.METADATA, order=order,
                                text=notes.strip(), value={"kind": "speaker_notes"},
                                slide=s_index + 1, extraction_confidence=1.0,
                            )
                        )
                        order += 1
            except Exception:  # noqa: BLE001
                pass

        return NirDocument(
            source_id=source_id,
            media_kind=MediaKind.SLIDES.value,
            version=version,
            engine=self._engine,
            engine_version=_pptx_version(),
            elements=tuple(elements),
            images=tuple(images),
            slides=slide_count,
            normalized_text="\n".join(text_parts)[: 8_000_000],
            needs_ocr=False,
        )


def _pptx_version() -> int:
    try:
        import pptx

        return int(pptx.__version__.split(".")[0])
    except Exception:  # noqa: BLE001
        return 0
